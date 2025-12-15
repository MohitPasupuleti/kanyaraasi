"""
Extractors Module - Secure Multi-Modal File Processing

This module handles SECURITY-FIRST file content extraction and processing.
CRITICAL: All files are processed locally BEFORE any external API calls to
prevent data leakage.

Architecture:
1. Extract text locally (NO external API calls)
2. Run LOCAL regex secrets detection
3. If secrets found → BLOCK immediately, Gemini NEVER called
4. Run LOCAL regex PII detection
5. If PII found → Sanitize locally FIRST
6. Only send sanitized/clean content to Gemini for context analysis

Supported file types:
- Images (.png, .jpg, .jpeg, .webp): OCR with pytesseract
- PDFs (.pdf): Text extraction with PyPDF2
- Word Documents (.docx): python-docx
- Excel Spreadsheets (.xlsx): openpyxl
- PowerPoint (.pptx): python-pptx
"""

import os
import re
from typing import Dict, List, Optional, Tuple
from pathlib import Path

# Image processing
try:
    import pytesseract
    from PIL import Image

    # Configure Tesseract path for Windows
    import platform
    if platform.system() == "Windows":
        # Try common Windows installation paths
        import os as os_module
        tesseract_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        ]
        for path in tesseract_paths:
            if os_module.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break

    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False

# PDF processing
try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

# Office document processing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import existing detection and sanitization functions
from executor import (
    detect_secrets,
    detect_pii,
    sanitize_with_gemini,
    calculate_risk_score,
    generate_explanation,
    evaluate_policy,
    regex_sanitize
)
from memory import log_decision
import uuid
from datetime import datetime


def extract_text_from_image_with_gemini(image_path: str) -> Dict[str, any]:
    """
    Extract text from image using Gemini multimodal API (CLOUD-BASED).

    This is the PRODUCTION-READY method that works everywhere:
    - Local development
    - Docker containers
    - Cloud deployments (AWS, GCP, Azure, etc.)

    No local OCR installation required!

    Args:
        image_path (str): Path to the image file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    try:
        import google.generativeai as genai
        import os
        from PIL import Image as PILImage

        # Check if Gemini API key is available
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            return {
                "success": False,
                "text": "",
                "error": "GEMINI_API_KEY not set. Cannot extract text from image."
            }

        # Configure Gemini
        genai.configure(api_key=api_key)

        # Open image
        image = PILImage.open(image_path)

        # Try multiple models with fallback (to handle rate limits)
        models_to_try = [
            'gemini-2.5-flash-lite',      # Has quota available
            'gemini-1.5-flash',            # Backup option
            'gemini-1.5-flash-8b',         # Another backup
            'gemini-2.0-flash-exp',        # Original (might be rate limited)
        ]

        prompt = """Extract ALL text from this image. Return ONLY the extracted text, nothing else.
If there is no text in the image, respond with: NO_TEXT_FOUND"""

        last_error = None
        for model_name in models_to_try:
            try:
                print(f"[OCR] Trying Gemini model: {model_name}")
                model = genai.GenerativeModel(model_name)
                response = model.generate_content([prompt, image])
                extracted_text = response.text.strip()
                print(f"[OCR] Successfully used model: {model_name}")
                break  # Success! Exit the loop
            except Exception as e:
                last_error = e
                print(f"[OCR] Model {model_name} failed: {str(e)}, trying next...")
                continue
        else:
            # All models failed
            raise Exception(f"All Gemini models failed. Last error: {last_error}")

        if extracted_text == "NO_TEXT_FOUND" or not extracted_text:
            return {
                "success": True,
                "text": "",
                "warning": "No text detected in image",
                "method": "gemini_vision"
            }

        return {
            "success": True,
            "text": extracted_text,
            "method": "gemini_vision"
        }

    except Exception as e:
        return {
            "success": False,
            "text": "",
            "error": f"Gemini Vision OCR failed: {str(e)}"
        }


def extract_text_locally_from_image(image_path: str) -> Dict[str, any]:
    """
    Extract text from image using LOCAL OCR ONLY (SECURITY-FIRST).

    This maintains the zero-leak architecture:
    1. Extract text LOCALLY using Tesseract (NO external API calls)
    2. If Tesseract unavailable, FAIL SECURE (don't send to Gemini)
    3. Let the caller decide whether to use Gemini Vision after secrets check

    Args:
        image_path (str): Path to the image file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    # SECURITY: Only use local Tesseract OCR
    # We do NOT send raw images to external APIs before checking for secrets
    if PYTESSERACT_AVAILABLE:
        try:
            from PIL import Image
            # Open and process image
            image = Image.open(image_path)

            # Extract text using OCR
            text = pytesseract.image_to_string(image)

            if not text or not text.strip():
                return {
                    "success": True,
                    "text": "",
                    "warning": "No text detected in image",
                    "method": "tesseract_local"
                }

            return {
                "success": True,
                "text": text.strip(),
                "method": "tesseract_local"
            }

        except Exception as e:
            # SECURITY: Fail secure - don't send to Gemini without secrets check
            return {
                "success": False,
                "text": "",
                "error": f"Local OCR failed: {str(e)}. Install Tesseract for secure image processing.",
                "method": "local_only"
            }

    # SECURITY: Tesseract not available - fail secure
    # We do NOT automatically fall back to Gemini Vision
    return {
        "success": False,
        "text": "",
        "error": "Tesseract OCR not installed. Please install Tesseract for secure local image processing. See: https://github.com/UB-Mannheim/tesseract/wiki",
        "method": "local_only"
    }


def extract_text_locally_from_pdf(pdf_path: str) -> Dict[str, any]:
    """
    Extract text from PDF using PyPDF2 (LOCAL, no API calls).

    Args:
        pdf_path (str): Path to the PDF file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    if not PYPDF2_AVAILABLE:
        return {
            "success": False,
            "text": "",
            "error": "PyPDF2 not installed. Run: pip install PyPDF2"
        }

    try:
        reader = PdfReader(pdf_path)

        # Extract text from all pages
        text_parts = []
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {page_num + 1}]\n{page_text}")

        full_text = "\n\n".join(text_parts)

        if not full_text.strip():
            return {
                "success": True,
                "text": "",
                "warning": "No text detected in PDF (might be image-based)"
            }

        return {
            "success": True,
            "text": full_text.strip()
        }

    except Exception as e:
        return {
            "success": False,
            "text": "",
            "error": f"PDF extraction failed: {str(e)}"
        }


def extract_text_locally_from_docx(docx_path: str) -> Dict[str, any]:
    """
    Extract text from Word document using python-docx (LOCAL, no API calls).

    Args:
        docx_path (str): Path to the .docx file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    if not DOCX_AVAILABLE:
        return {
            "success": False,
            "text": "",
            "error": "python-docx not installed. Run: pip install python-docx"
        }

    try:
        doc = Document(docx_path)

        # Extract text from paragraphs
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

        # Extract text from tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                if row_text.strip():
                    table_texts.append(row_text)

        # Combine all text
        all_text = "\n".join(paragraphs)
        if table_texts:
            all_text += "\n\n[Tables]\n" + "\n".join(table_texts)

        if not all_text.strip():
            return {
                "success": True,
                "text": "",
                "warning": "No text detected in document"
            }

        return {
            "success": True,
            "text": all_text.strip()
        }

    except Exception as e:
        return {
            "success": False,
            "text": "",
            "error": f"DOCX extraction failed: {str(e)}"
        }


def extract_text_locally_from_xlsx(xlsx_path: str) -> Dict[str, any]:
    """
    Extract text from Excel spreadsheet using openpyxl (LOCAL, no API calls).

    Args:
        xlsx_path (str): Path to the .xlsx file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    if not OPENPYXL_AVAILABLE:
        return {
            "success": False,
            "text": "",
            "error": "openpyxl not installed. Run: pip install openpyxl"
        }

    try:
        workbook = load_workbook(xlsx_path, read_only=True, data_only=True)

        all_text = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            all_text.append(f"[Sheet: {sheet_name}]")

            # Extract cell values
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    all_text.append(row_text)

        full_text = "\n".join(all_text)

        if not full_text.strip():
            return {
                "success": True,
                "text": "",
                "warning": "No text detected in spreadsheet"
            }

        return {
            "success": True,
            "text": full_text.strip()
        }

    except Exception as e:
        return {
            "success": False,
            "text": "",
            "error": f"XLSX extraction failed: {str(e)}"
        }


def extract_text_locally_from_pptx(pptx_path: str) -> Dict[str, any]:
    """
    Extract text from PowerPoint using python-pptx (LOCAL, no API calls).

    Args:
        pptx_path (str): Path to the .pptx file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    if not PPTX_AVAILABLE:
        return {
            "success": False,
            "text": "",
            "error": "python-pptx not installed. Run: pip install python-pptx"
        }

    try:
        prs = Presentation(pptx_path)

        all_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            all_text.append(f"[Slide {slide_num}]")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text)

        full_text = "\n".join(all_text)

        if not full_text.strip():
            return {
                "success": True,
                "text": "",
                "warning": "No text detected in presentation"
            }

        return {
            "success": True,
            "text": full_text.strip()
        }

    except Exception as e:
        return {
            "success": False,
            "text": "",
            "error": f"PPTX extraction failed: {str(e)}"
        }


def extract_text_locally_from_txt(txt_path: str) -> Dict[str, any]:
    """
    Extract text from plain text file (LOCAL, no API calls).

    Args:
        txt_path (str): Path to the .txt file

    Returns:
        dict: Extraction result with 'success', 'text', and optional 'error'
    """
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            text = f.read()

        if not text or not text.strip():
            return {
                "success": True,
                "text": "",
                "warning": "No text detected in file"
            }

        return {
            "success": True,
            "text": text.strip()
        }

    except Exception as e:
        return {
            "success": False,
            "text": "",
            "error": f"TXT extraction failed: {str(e)}"
        }


def get_file_metadata(file_path: str) -> Dict[str, any]:
    """
    Extract metadata from a file.

    Args:
        file_path (str): Path to the file

    Returns:
        dict: File metadata including name, size, type, etc.
    """
    try:
        path = Path(file_path)

        return {
            "filename": path.name,
            "extension": path.suffix.lower(),
            "size_bytes": path.stat().st_size,
            "size_human": _human_readable_size(path.stat().st_size),
            "exists": path.exists(),
            "is_file": path.is_file()
        }
    except Exception as e:
        return {
            "filename": os.path.basename(file_path),
            "extension": os.path.splitext(file_path)[1].lower(),
            "error": str(e)
        }


def _human_readable_size(size_bytes: int) -> str:
    """Convert bytes to human-readable format."""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.2f} TB"


def sanitize_text_locally(text: str, pii_patterns: List[str]) -> str:
    """
    Sanitize text locally using regex patterns BEFORE any API call.

    This ensures PII is removed locally before sending to external APIs.

    Args:
        text (str): Text to sanitize
        pii_patterns (list): List of PII types found (e.g., ['email', 'phone'])

    Returns:
        str: Sanitized text with PII replaced by [REDACTED] markers
    """
    return regex_sanitize(text)


def process_file_securely(file_path: str, file_type: str, use_case: str = "general") -> Dict[str, any]:
    """
    MAIN SECURE PROCESSING FUNCTION - Security-First File Analysis

    CRITICAL SECURITY FLOW:
    1. Extract text locally (no API calls)
    2. Run LOCAL regex secrets detection
    3. If secrets found → BLOCK immediately, Gemini NEVER called
    4. Run LOCAL regex PII detection
    5. If PII found → Sanitize locally FIRST
    6. Only send sanitized/clean content to Gemini for context analysis

    Args:
        file_path (str): Path to the file to process
        file_type (str): File type/extension (e.g., '.png', '.pdf', '.docx')
        use_case (str): Context of request (debugging, support, docs, general)

    Returns:
        dict: Decision result matching the format of process_input() from planner.py
    """
    # Initialize execution context
    execution_trace = []
    audit_id = str(uuid.uuid4())
    timestamp = datetime.utcnow().isoformat() + "Z"
    gemini_called = False
    processing_method = "local_only"

    # Get file metadata
    file_info = get_file_metadata(file_path)
    execution_trace.append(f"OBSERVE: Processing file {file_info.get('filename', 'unknown')}")
    execution_trace.append(f"   File size: {file_info.get('size_human', 'unknown')}")
    execution_trace.append(f"   File type: {file_type}")
    execution_trace.append(f"   Use case: {use_case}")

    # PHASE 1: LOCAL TEXT EXTRACTION (NO API CALLS)
    execution_trace.append("ACT: Phase 1 - LOCAL text extraction (no API calls)")

    # Determine extraction method based on file type
    file_type_lower = file_type.lower()

    if file_type_lower in ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif']:
        extraction_result = extract_text_locally_from_image(file_path)
        extraction_method = "pytesseract OCR"
    elif file_type_lower == '.pdf':
        extraction_result = extract_text_locally_from_pdf(file_path)
        extraction_method = "PyPDF2"
    elif file_type_lower == '.docx':
        extraction_result = extract_text_locally_from_docx(file_path)
        extraction_method = "python-docx"
    elif file_type_lower == '.xlsx':
        extraction_result = extract_text_locally_from_xlsx(file_path)
        extraction_method = "openpyxl"
    elif file_type_lower == '.pptx':
        extraction_result = extract_text_locally_from_pptx(file_path)
        extraction_method = "python-pptx"
    elif file_type_lower == '.txt':
        extraction_result = extract_text_locally_from_txt(file_path)
        extraction_method = "plain text"
    else:
        return {
            "decision": "BLOCK",
            "risk_score": 50,
            "explanation": f"Unsupported file type: {file_type}. Only images, PDFs, and Office documents are supported.",
            "file_info": file_info,
            "gemini_called": False,
            "processing_method": "unsupported",
            "audit_id": audit_id,
            "timestamp": timestamp,
            "detected_signals": [],
            "policy_refs": []
        }

    # Check extraction success
    if not extraction_result["success"]:
        error_msg = extraction_result.get("error", "Unknown extraction error")
        execution_trace.append(f"   ERROR: Local extraction failed - {error_msg}")

        # SECURITY OPTION: Allow Gemini Vision as fallback for images ONLY
        # This is controlled by environment variable for security
        use_gemini_fallback = os.getenv("USE_GEMINI_VISION_FOR_OCR", "false").lower() == "true"

        if file_type_lower in ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif'] and use_gemini_fallback:
            execution_trace.append("   FALLBACK: Attempting Gemini Vision OCR (environment variable enabled)")
            print("[SECURITY WARNING] Using Gemini Vision API as fallback - raw image will be sent to Google")
            print("                   To disable, set USE_GEMINI_VISION_FOR_OCR=false in .env")

            gemini_extraction = extract_text_from_image_with_gemini(file_path)
            if gemini_extraction["success"]:
                extracted_text = gemini_extraction["text"]
                extraction_method = "Gemini Vision (fallback)"
                execution_trace.append(f"   Gemini Vision extraction successful: {len(extracted_text)} characters")
            else:
                return {
                    "decision": "BLOCK",
                    "risk_score": 50,
                    "explanation": f"Both local OCR and Gemini Vision failed: {gemini_extraction.get('error', 'Unknown error')}",
                    "file_info": file_info,
                    "gemini_called": False,
                    "processing_method": "extraction_failed",
                    "audit_id": audit_id,
                    "timestamp": timestamp,
                    "detected_signals": [],
                    "policy_refs": []
                }
        else:
            # FAIL SECURE: Don't process images without local OCR
            return {
                "decision": "BLOCK",
                "risk_score": 50,
                "explanation": f"File processing failed: {error_msg}. Install Tesseract for local image processing or enable USE_GEMINI_VISION_FOR_OCR=true (not recommended for sensitive content).",
                "file_info": file_info,
                "gemini_called": False,
                "processing_method": "extraction_failed",
                "audit_id": audit_id,
                "timestamp": timestamp,
                "detected_signals": [],
                "policy_refs": []
            }
    else:
        extracted_text = extraction_result["text"]
        execution_trace.append(f"   Successfully extracted {len(extracted_text)} characters using {extraction_method}")

    # Debug: Print extracted text to see what OCR found
    print(f"[OCR] Extracted text preview (first 500 chars):")
    print(f"       {extracted_text[:500] if extracted_text else '(empty)'}")

    if extraction_result.get("warning"):
        execution_trace.append(f"   Warning: {extraction_result['warning']}")

    # Handle empty content
    if not extracted_text or not extracted_text.strip():
        execution_trace.append("   No text content found in file")
        execution_trace.append("DECIDE: ALLOW (empty file, no sensitive data)")

        log_decision("ALLOW", "No text content detected in file", execution_trace, audit_id)

        return {
            "decision": "ALLOW",
            "risk_score": 0,
            "explanation": "No text content detected in file. File appears to be empty or contain only images.",
            "file_info": file_info,
            "gemini_called": False,
            "processing_method": "local_only",
            "audit_id": audit_id,
            "timestamp": timestamp,
            "detected_signals": [{"type": "Empty", "severity": "none", "location": "N/A", "description": "No text detected"}],
            "policy_refs": []
        }

    # PHASE 2: LOCAL SECRETS DETECTION (NO API CALLS)
    execution_trace.append("ACT: Phase 2 - LOCAL secrets detection (no API calls)")
    secrets_result = detect_secrets(extracted_text)

    if secrets_result["found"]:
        # CRITICAL: EARLY EXIT - Secrets found, BLOCK immediately
        execution_trace.append(f"   CRITICAL: Secret detected: {secrets_result['pattern']}")
        execution_trace.append("EARLY EXIT: Blocking due to secrets (Gemini NEVER called)")
        execution_trace.append("   Zero-leak guarantee: Secret never left your infrastructure ✓")

        policy_result = evaluate_policy(use_case, True, False, [])
        risk_score = calculate_risk_score(True, False, [], False)

        explanation = generate_explanation(
            "BLOCK",
            risk_score,
            True,
            False,
            [],
            policy_result["policy_refs"]
        )

        detected_signals = [{
            "type": secrets_result["pattern"],
            "severity": "critical",
            "location": f"Found: {secrets_result.get('matched_value', 'N/A')}",
            "description": f"Detected {secrets_result['pattern']} in file content"
        }]

        log_decision("BLOCK", explanation, execution_trace, audit_id)

        return {
            "decision": "BLOCK",
            "risk_score": risk_score,
            "explanation": explanation,
            "policy_refs": policy_result["policy_refs"],
            "detected_signals": detected_signals,
            "file_info": file_info,
            "gemini_called": False,
            "processing_method": "local_only",
            "audit_id": audit_id,
            "timestamp": timestamp
        }

    execution_trace.append("   No secrets detected ✓")

    # PHASE 3: LOCAL PII DETECTION (NO API CALLS)
    execution_trace.append("ACT: Phase 3 - LOCAL PII detection (no API calls)")
    pii_result = detect_pii(extracted_text)

    pii_found = pii_result["found"]
    pii_types = pii_result["types"]

    if pii_found:
        execution_trace.append(f"   PII detected: {', '.join(pii_types)}")
    else:
        execution_trace.append("   No PII detected ✓")

    # PHASE 4: LOCAL SANITIZATION (if PII found)
    sanitized_text = extracted_text
    sanitization_applied = False

    if pii_found:
        execution_trace.append("ACT: Phase 4 - LOCAL sanitization (before any API call)")
        sanitized_text = sanitize_text_locally(extracted_text, pii_types)
        sanitization_applied = True
        execution_trace.append("   PII sanitized locally ✓")

    # PHASE 5: POLICY EVALUATION
    execution_trace.append("ACT: Phase 5 - Policy evaluation")
    policy_result = evaluate_policy(use_case, False, pii_found, pii_types)
    execution_trace.append(f"   Policy: {policy_result['policy_refs'][0]['id']}")

    # PHASE 6: OPTIONAL GEMINI CALL (only with sanitized content)
    gemini_result = None

    if policy_result["allow_sanitization"] and sanitization_applied:
        execution_trace.append("ACT: Phase 6 - Gemini analysis (on SANITIZED content only)")
        gemini_result = sanitize_with_gemini(sanitized_text)
        gemini_called = True
        processing_method = "local_then_gemini"

        if gemini_result["used_llm"]:
            execution_trace.append(f"   Gemini analysis complete (model: {gemini_result.get('model_used', 'gemini')})")
            sanitized_text = gemini_result["sanitized_text"]
        else:
            execution_trace.append("   Gemini unavailable, using local sanitization")

    # PHASE 7: RISK SCORING
    execution_trace.append("CALCULATE: Computing risk score")
    risk_score = calculate_risk_score(False, pii_found, pii_types, sanitization_applied)
    execution_trace.append(f"   Risk score: {risk_score}/100")

    # PHASE 8: EXPLANATION
    execution_trace.append("EXPLAIN: Generating decision explanation")
    explanation = generate_explanation(
        "SANITIZE" if sanitization_applied else "ALLOW",
        risk_score,
        False,
        pii_found,
        pii_types,
        policy_result["policy_refs"]
    )

    # PHASE 9: DECISION
    decision = "SANITIZE" if sanitization_applied else "ALLOW"
    execution_trace.append(f"DECIDE: {decision}")

    # Build detected signals
    detected_signals = []
    if pii_found:
        detected_signals.append({
            "type": f"PII: {', '.join(pii_types)}",
            "severity": "medium",
            "location": "File content",
            "description": f"Detected {len(pii_types)} type(s) of PII"
        })
    if sanitization_applied:
        detected_signals.append({
            "type": "Sanitized",
            "severity": "low",
            "location": "Full content",
            "description": f"Sanitized {'with Gemini' if gemini_result and gemini_result['used_llm'] else 'locally'}"
        })

    if not detected_signals:
        detected_signals.append({
            "type": "Clean",
            "severity": "none",
            "location": "N/A",
            "description": "No issues detected"
        })

    # LOG DECISION
    log_decision(decision, explanation, execution_trace, audit_id)

    # Build response
    response = {
        "decision": decision,
        "risk_score": risk_score,
        "explanation": explanation,
        "policy_refs": policy_result["policy_refs"],
        "detected_signals": detected_signals,
        "original_content": extracted_text,  # ADD: Show extracted text in frontend
        "file_info": file_info,
        "gemini_called": gemini_called,
        "processing_method": processing_method,
        "audit_id": audit_id,
        "timestamp": timestamp
    }

    # Add sanitized content if available
    if sanitization_applied:
        response["safe_prompt"] = sanitized_text  # ADD: For frontend copy button
        response["safe_content"] = sanitized_text
        response["masked_content"] = sanitized_text
        response["diff"] = {
            "original_length": len(extracted_text),
            "sanitized_length": len(sanitized_text),
            "reduction": len(extracted_text) - len(sanitized_text)
        }
    else:
        # Even if not sanitized, provide safe content (same as original)
        response["safe_prompt"] = extracted_text
        response["masked_content"] = extracted_text

    return response
